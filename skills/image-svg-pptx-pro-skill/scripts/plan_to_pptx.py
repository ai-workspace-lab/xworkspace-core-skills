#!/usr/bin/env python3
"""Generate an editable PPTX from layout_plan.json.

This is the preferred converter because the plan preserves semantics better than
arbitrary SVG. It maps supported plan elements to native PowerPoint objects and
uses image assets for complex regions.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def hex_to_rgb(value: str, default: Tuple[int, int, int] = (0, 0, 0)) -> RGBColor:
    if not value or value == "none":
        return RGBColor(*default)
    value = str(value).strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*default)


class Mapper:
    def __init__(self, canvas_w: float, canvas_h: float, slide_w_in: float = 13.333333):
        self.canvas_w = canvas_w
        self.canvas_h = canvas_h
        self.slide_w = Inches(slide_w_in)
        self.slide_h = int(self.slide_w * canvas_h / canvas_w)
        self.pt_per_px = (slide_w_in * 72.0) / canvas_w

    def x(self, px: float):
        return int(self.slide_w * float(px) / self.canvas_w)

    def y(self, px: float):
        return int(self.slide_h * float(px) / self.canvas_h)

    def w(self, px: float):
        return int(self.slide_w * float(px) / self.canvas_w)

    def h(self, px: float):
        return int(self.slide_h * float(px) / self.canvas_h)

    def pt(self, px: float):
        # The factor makes image-pixel typography visually close on standard 16:9 decks.
        return Pt(max(1, float(px) * self.pt_per_px))


def apply_fill(shape, fill_value: str | None):
    if fill_value is None or fill_value == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_value, (255, 255, 255))


def apply_line(shape, stroke: str | None, stroke_width_pt: float | None):
    if stroke is None or stroke == "none":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(stroke)
        if stroke_width_pt is not None:
            shape.line.width = Pt(max(0.25, stroke_width_pt))


def add_rect(slide, m: Mapper, el: Dict[str, Any]):
    x, y, w, h = m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 0)), m.h(el.get("h", 0))
    radius = float(el.get("rx", 0) or 0)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    apply_fill(shape, el.get("fill", "#FFFFFF"))
    apply_line(shape, el.get("stroke", "none"), float(el.get("stroke_width", 0) or 0) * m.pt_per_px)
    return shape


def add_line(slide, m: Mapper, el: Dict[str, Any]):
    shape = slide.shapes.add_connector(
        1,
        m.x(el.get("x1", 0)),
        m.y(el.get("y1", 0)),
        m.x(el.get("x2", 0)),
        m.y(el.get("y2", 0)),
    )
    apply_line(shape, el.get("stroke", "#000000"), float(el.get("stroke_width", 1) or 1) * m.pt_per_px)
    return shape


def add_circle(slide, m: Mapper, el: Dict[str, Any]):
    cx = float(el.get("cx", el.get("x", 0)))
    cy = float(el.get("cy", el.get("y", 0)))
    r = float(el.get("r", 0))
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, m.x(cx - r), m.y(cy - r), m.w(2 * r), m.h(2 * r))
    apply_fill(shape, el.get("fill", "#FFFFFF"))
    apply_line(shape, el.get("stroke", "none"), float(el.get("stroke_width", 0) or 0) * m.pt_per_px)
    return shape


def add_text(slide, m: Mapper, el: Dict[str, Any]):
    box = slide.shapes.add_textbox(m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    valign = str(el.get("valign", "top"))
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    lines = str(el.get("text", "")).split("\n")
    p = tf.paragraphs[0]
    p.text = lines[0] if lines else ""
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.text = extra
    align = str(el.get("align", "left"))
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    for p in tf.paragraphs:
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        for run in p.runs:
            run.font.name = el.get("font_family", "Microsoft YaHei")
            run.font.size = m.pt(float(el.get("font_size", 24)))
            run.font.bold = int(el.get("font_weight", 400) or 400) >= 600
            run.font.color.rgb = hex_to_rgb(el.get("color", "#000000"))
    return box


def add_image(slide, m: Mapper, el: Dict[str, Any], assets_dir: Path, asset_map: Dict[str, Any]):
    asset = asset_map.get(el.get("asset_id"), {})
    file = asset.get("file") or el.get("file")
    if not file:
        return None
    path = assets_dir / file
    if not path.exists():
        # Also allow absolute or plan-relative file paths.
        path = Path(file)
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 0)), m.h(el.get("h", 0)))


def add_table(slide, m: Mapper, el: Dict[str, Any]):
    rows = int(el.get("rows", 1))
    cols = int(el.get("cols", 1))
    shape = slide.shapes.add_table(rows, cols, m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    table = shape.table
    cell_text = el.get("cell_text", [])
    for r in range(rows):
        for c in range(cols):
            text = ""
            if r < len(cell_text) and c < len(cell_text[r]):
                text = str(cell_text[r][c])
            cell = table.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = el.get("font_family", "Microsoft YaHei")
                    run.font.size = m.pt(float(el.get("font_size", 18)))
                    run.font.color.rgb = hex_to_rgb(el.get("color", "#111111"))
    return shape


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path)
    parser.add_argument("--assets", type=Path, default=Path("work/assets"))
    parser.add_argument("--out", type=Path, default=Path("work/reconstructed.pptx"))
    parser.add_argument("--slide-width", type=float, default=13.333333)
    args = parser.parse_args()

    plan = json.loads(args.plan.read_text(encoding="utf-8"))
    canvas = plan["canvas"]
    m = Mapper(float(canvas["width"]), float(canvas["height"]), args.slide_width)

    prs = Presentation()
    prs.slide_width = m.slide_w
    prs.slide_height = m.slide_h
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background as native filled rectangle.
    bg_color = canvas.get("background", "#FFFFFF")
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    apply_fill(bg, bg_color)
    apply_line(bg, "none", None)

    asset_map = {a.get("id"): a for a in plan.get("assets", [])}
    elements = sorted(plan.get("elements", []), key=lambda e: e.get("z", 0))

    for el in elements:
        t = el.get("type")
        if t == "rect":
            add_rect(slide, m, el)
        elif t == "line":
            add_line(slide, m, el)
        elif t == "circle":
            add_circle(slide, m, el)
        elif t == "text":
            add_text(slide, m, el)
        elif t == "image":
            add_image(slide, m, el, args.assets, asset_map)
        elif t == "table":
            add_table(slide, m, el)
        else:
            # Unsupported elements should be present as cropped assets/fallbacks.
            continue

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(str(args.out))


if __name__ == "__main__":
    main()
